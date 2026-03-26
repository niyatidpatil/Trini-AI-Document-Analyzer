# core/loader.py
from __future__ import annotations
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Dict, Optional
import unicodedata

import io
import requests
from bs4 import BeautifulSoup

# PDF – high fidelity text
import fitz  # PyMuPDF
print(fitz.__doc__)
# OCR fallback (optional; guarded)
try:
    from PIL import Image
    import pytesseract
    _OCR_AVAILABLE = True
except Exception:
    _OCR_AVAILABLE = False

# Office formats (optional; guarded per format)
try:
    import docx  # python-docx
    _DOCX_AVAILABLE = True
except Exception:
    _DOCX_AVAILABLE = False

try:
    import pptx  # python-pptx
    _PPTX_AVAILABLE = True
except Exception:
    _PPTX_AVAILABLE = False

try:
    import pandas as pd  # for xlsx via openpyxl engine
    _XLSX_AVAILABLE = True
except Exception:
    _XLSX_AVAILABLE = False


@dataclass
class Doc:
    content: str
    metadata: dict

def _normalize_text_basic(text: str) -> str:
    """Gentle cleanup: unicode normalize, remove soft hyphens/tabs/CR,
    de-hyphenate across line breaks, collapse single LFs to spaces but
    preserve double LFs as paragraph breaks.
    """
    import re, unicodedata
    t = unicodedata.normalize("NFKC", text or "")
    t = t.replace("\r", "").replace("\t", " ").replace("\u00ad", "")
    t = re.sub(r"(\w)-\n(\w)", r"\1\2", t)  # de-hyphenate across LB
    t = t.replace("\n\n", "<<<PARA>>>")
    t = re.sub(r"[ \t]*\n[ \t]*", " ", t)  # single LB -> space
    t = re.sub(r"[ ]{2,}", " ", t)         # collapse spaces
    t = t.replace("<<<PARA>>>", "\n\n")
    return t.strip()


def _detect_repeating_lines(
    pages_lines: List[List[str]],
    top_k: int = 2,
    bottom_k: int = 2,
    min_fraction: float = 0.30,
) -> Dict[str, List[str]]:
    """Heuristic: find lines that repeat at the top/bottom of many pages.
    Returns lists of header/footer patterns to strip later.
    """
    from collections import Counter

    top_counter, bot_counter = Counter(), Counter()
    n = len(pages_lines) or 1

    for lines in pages_lines:
        if not lines:
            continue
        # sample top K non-empty lines
        top = [ln.strip() for ln in lines[:top_k] if ln.strip()]
        # sample bottom K non-empty lines
        bot = [ln.strip() for ln in lines[-bottom_k:] if ln.strip()]
        top_counter.update(top)
        bot_counter.update(bot)

    threshold = max(2, int(min_fraction * n))  # appear on >=30% of pages
    header = [s for s, c in top_counter.items() if c >= threshold]
    footer = [s for s, c in bot_counter.items() if c >= threshold]

    return {"header": header, "footer": footer}


# ---------- Extraction reporting (for transparency in UI) ----------
class ExtractionReport:
    def __init__(self) -> None:
        # Totals
        self.pdf_pages_total = 0
        self.pdf_pages_text_ok = 0
        self.pdf_pages_text_weak = 0
        self.pdf_pages_ocr_used = 0
        self.pdf_pages_ocr_failed = 0
        self.pdf_pages_errors: List[int] = []

        # Per-page details
        # each item: {
        #   "page": int,
        #   "chars": int,
        #   "weak": bool,
        #   "ocr_used": bool,
        #   "header_removed": int,
        #   "footer_removed": int,
        #   "error": Optional[str],
        # }
        self.pages: List[Dict] = []

        # Detected repeating patterns (sampled from lines)
        self.header_patterns: List[str] = []
        self.footer_patterns: List[str] = []

        # Free-form notes
        self.messages: List[str] = []

    def to_dict(self) -> Dict:
        return {
            "pdf_pages_total": self.pdf_pages_total,
            "pdf_pages_text_ok": self.pdf_pages_text_ok,
            "pdf_pages_text_weak": self.pdf_pages_text_weak,
            "pdf_pages_ocr_used": self.pdf_pages_ocr_used,
            "pdf_pages_ocr_failed": self.pdf_pages_ocr_failed,
            "pdf_pages_errors": self.pdf_pages_errors,
            "header_patterns": self.header_patterns,
            "footer_patterns": self.footer_patterns,
            "pages": self.pages,
            "messages": self.messages,
        }


_last_report = ExtractionReport()


def _reset_report():
    global _last_report
    _last_report = ExtractionReport()


def get_last_extraction_report() -> Dict:
    """Return a summary of the most recent PDF extraction pass."""
    return _last_report.to_dict()

def _normalize_text(text: str) -> str:
    """
    Make extracted PDF text read smoothly:
    - Unicode normalize (NFKC)
    - Remove soft hyphens
    - De-hyphenate across line breaks: e.g., 'risk-\nadjusted' -> 'risk-adjusted'
    - Preserve paragraphs: keep double newlines, collapse single newlines to spaces
    - Collapse excess spaces
    """
    if not text:
        return ""

    # Unicode normalize (fixes ligatures, spacing, etc.)
    text = unicodedata.normalize("NFKC", text)

    # Remove soft hyphen
    text = text.replace("\u00ad", "")

    # Standardize line endings and tabs
    text = text.replace("\r", "")
    text = text.replace("\t", " ")

    # De-hyphenate where a hyphen splits a word across a newline
    # (only letters on both sides to avoid math/code side-effects)
    import re
    text = re.sub(r"([A-Za-z])-\n([A-Za-z])", r"\1\2", text)

    # Preserve paragraphs: protect double newlines first
    PARA = "<<<PARA>>>"
    text = text.replace("\n\n", PARA)

    # Now collapse single newlines to spaces
    text = re.sub(r"[ \t]*\n[ \t]*", " ", text)

    # Restore paragraph breaks
    text = text.replace(PARA, "\n\n")

    # Collapse multiple spaces
    text = re.sub(r"[ ]{2,}", " ", text)

    return text.strip()

# ---------- PDF extraction with PyMuPDF + optional OCR fallback ----------
def _extract_pdf_text_with_pymupdf(
    path: Path,
    weak_threshold_chars: int = 60,   # pages below this count considered "weak"
    ocr_dpi: int = 200,
) -> str:
    """
    Extract text in reading order using PyMuPDF blocks.
    If a page has too few chars, optionally tries OCR (if Tesseract is available).
    Logs stats to _last_report.
    NOTE: This step does NOT remove headers/footers yet (that comes in Step 3).
    """
    global _last_report
    text_parts: List[str] = []

    try:
        doc = fitz.open(str(path))
    except Exception as e:
        _last_report.messages.append(f"[ERROR] Failed to open PDF: {path.name} ({e})")
        return ""

    _last_report.pdf_pages_total += len(doc)

    for page_index in range(len(doc)):
        try:
            page = doc.load_page(page_index)

            # --- Block-based extraction (more layout-aware than "text") ---
            # Each block: (x0, y0, x1, y1, text, block_no, ...)
            blocks = page.get_text("blocks") or []
            # Sort by vertical (y) then horizontal (x) to approximate reading order (handles 2 columns well)
            blocks.sort(key=lambda b: (round(b[1], 1), round(b[0], 1)))

            # Light filter: skip empty/whitespace-only blocks; keep everything else
            block_texts = []
            for b in blocks:
                txt = b[4] if len(b) > 4 else ""
                if txt and txt.strip():
                    block_texts.append(txt.strip())

            page_text = "\n".join(block_texts)
            page_text = _normalize_text(page_text)

            if len(page_text) >= weak_threshold_chars:
                _last_report.pdf_pages_text_ok += 1
                text_parts.append(page_text)
                continue

            # --- Weak text: try OCR fallback if available ---
            _last_report.pdf_pages_text_weak += 1
            if _OCR_AVAILABLE:
                try:
                    pix = page.get_pixmap(dpi=ocr_dpi, alpha=False)
                    img_bytes = pix.tobytes("png")
                    pil_img = Image.open(io.BytesIO(img_bytes))
                    ocr_text = pytesseract.image_to_string(pil_img)
                    ocr_text = _normalize_text(ocr_text)
                    if ocr_text and len(ocr_text) >= weak_threshold_chars:
                        _last_report.pdf_pages_ocr_used += 1
                        text_parts.append(ocr_text)
                    else:
                        _last_report.pdf_pages_ocr_failed += 1
                except Exception:
                    _last_report.pdf_pages_ocr_failed += 1
            # If OCR not available, we just skip adding weak text

        except Exception:
            _last_report.pdf_pages_errors.append(page_index)

    doc.close()
    return "\n".join(text_parts)


def _extract_pdf_pages_with_pymupdf(
    path: Path,
    weak_threshold_chars: int = 60,
    ocr_dpi: int = 200,
) -> List[str]:
    """
    Enhanced page-by-page extraction:
    - Layout-aware block order (Y then X)
    - Two-pass: detect/remove repeating headers/footers across pages
    - Gentle normalization (de-hyphenation, paragraph preservation)
    - Optional OCR fallback on weak pages
    - Rich per-page reporting
    """
    global _last_report
    pages_out: List[str] = []

    # ------------- Pass 1: collect raw per-page lines -------------
    raw_lines_per_page: List[List[str]] = []
    open_error = None
    try:
        doc = fitz.open(str(path))
    except Exception as e:
        open_error = str(e)
        _last_report.messages.append(f"[ERROR] Failed to open PDF: {path.name} ({e})")
        return pages_out

    _last_report.pdf_pages_total += len(doc)

    # We will also store preliminary stats to augment in Pass 2
    prelim_stats = [{"weak": False, "ocr_used": False, "error": None} for _ in range(len(doc))]

    for page_index in range(len(doc)):
        try:
            page = doc.load_page(page_index)
            blocks = page.get_text("blocks") or []
            blocks.sort(key=lambda b: (round(b[1], 1), round(b[0], 1)))  # (y, x)

            # Build raw page text from blocks, quickly split into lines
            raw_parts = []
            for b in blocks:
                if len(b) >= 5:
                    txt = (b[4] or "").strip()
                    if txt:
                        raw_parts.append(txt)

            raw_text = "\n".join(raw_parts).strip()

            if len(raw_text) < weak_threshold_chars:
                # Try OCR if available
                _last_report.pdf_pages_text_weak += 1
                prelim_stats[page_index]["weak"] = True

                if _OCR_AVAILABLE:
                    try:
                        pix = page.get_pixmap(dpi=ocr_dpi, alpha=False)
                        img_bytes = pix.tobytes("png")
                        pil_img = Image.open(io.BytesIO(img_bytes))
                        ocr_text = (pytesseract.image_to_string(pil_img) or "").strip()
                        if len(ocr_text) >= weak_threshold_chars:
                            _last_report.pdf_pages_ocr_used += 1
                            prelim_stats[page_index]["ocr_used"] = True
                            raw_text = ocr_text
                        else:
                            _last_report.pdf_pages_ocr_failed += 1
                    except Exception:
                        _last_report.pdf_pages_ocr_failed += 1
                # if OCR not available, keep raw_text (likely empty/weak)
            else:
                _last_report.pdf_pages_text_ok += 1

            # Split into trimmed lines now; we’ll strip header/footer later
            lines = [ln.strip() for ln in (raw_text.split("\n") if raw_text else [])]
            raw_lines_per_page.append(lines)

        except Exception as e:
            _last_report.pdf_pages_errors.append(page_index)
            prelim_stats[page_index]["error"] = str(e)
            raw_lines_per_page.append([])

    # ------------- Pass 2: detect and strip repeated header/footer -------------
    patterns = _detect_repeating_lines(
        raw_lines_per_page,
        top_k=2,         # look at up to 2 lines from top
        bottom_k=2,      # and 2 lines from bottom
        min_fraction=0.30
    )
    header_patterns = patterns.get("header", [])
    footer_patterns = patterns.get("footer", [])

    _last_report.header_patterns = header_patterns
    _last_report.footer_patterns = footer_patterns

    # Strip patterns and normalize page text; record per-page stats
    for i, lines in enumerate(raw_lines_per_page):
        header_removed = 0
        footer_removed = 0

        # remove matching headers from the *very top* only
        while lines and lines[0] in header_patterns:
            lines.pop(0)
            header_removed += 1

        # remove matching footers from the *very bottom* only
        while lines and lines[-1] in footer_patterns:
            lines.pop()
            footer_removed += 1

        cleaned = "\n".join(lines)
        cleaned = _normalize_text_basic(cleaned)
        pages_out.append(cleaned)

        # Per-page chars after cleanup
        char_count = len(cleaned)

        _last_report.pages.append({
            "page": i,
            "chars": char_count,
            "weak": bool(prelim_stats[i]["weak"]),
            "ocr_used": bool(prelim_stats[i]["ocr_used"]),
            "header_removed": header_removed,
            "footer_removed": footer_removed,
            "error": prelim_stats[i]["error"],
        })

    doc.close()
    return pages_out


# ---------- Simple TXT ----------
def _read_txt(path: Path) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


# ---------- Public loaders ----------
def load_pdfs(paths: Iterable[str | Path]) -> List[Doc]:
    """Load PDFs as one Doc per page, including page metadata for traceability."""
    _reset_report()
    docs: List[Doc] = []
    for p in paths:
        p = Path(p)
        if not p.exists():
            _last_report.messages.append(f"[WARN] File not found: {p}")
            continue

        page_texts = _extract_pdf_pages_with_pymupdf(p)
        total_pages = len(page_texts)
        for i, txt in enumerate(page_texts):
            # Even if txt is empty (weak page without OCR), we still emit a page Doc for transparency.
            meta = {
                "source": str(p),
                "file_name": p.name,
                "type": "pdf",
                "page": i,                 # 0-based; display as i+1 if you prefer
                "page_count": total_pages,
            }
            docs.append(Doc(content=txt or "", metadata=meta))
    return docs


def load_txts(paths: Iterable[str | Path]) -> List[Doc]:
    docs: List[Doc] = []
    for p in paths:
        p = Path(p)
        if not p.exists():
            continue
        content = _read_txt(p)
        docs.append(Doc(content=content, metadata={"source": str(p), "type": "txt"}))
    return docs


def load_urls(urls: Iterable[str]) -> List[Doc]:
    docs: List[Doc] = []
    for url in urls:
        u = (url or "").strip()
        if not u:
            continue
        try:
            r = requests.get(u, timeout=20)
            r.raise_for_status()
            soup = BeautifulSoup(r.text, "html.parser")
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()
            text = " ".join(soup.get_text(separator=" ").split())
            docs.append(Doc(content=text, metadata={"source": u, "type": "url"}))
        except Exception:
            continue
    return docs


# ---------- Office loaders (optional) ----------
def load_docx(paths: Iterable[str | Path]) -> List[Doc]:
    docs: List[Doc] = []
    if not _DOCX_AVAILABLE:
        return docs
    for p in paths:
        p = Path(p)
        if not p.exists():
            continue
        try:
            d = docx.Document(str(p))
            parts = []
            # paragraphs
            for para in d.paragraphs:
                parts.append(para.text)
            # simple tables
            for tbl in d.tables:
                for row in tbl.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
            content = "\n".join([t for t in parts if t and t.strip()])
            docs.append(Doc(content=content, metadata={"source": str(p), "type": "docx"}))
        except Exception:
            continue
    return docs


def load_pptx(paths: Iterable[str | Path]) -> List[Doc]:
    docs: List[Doc] = []
    if not _PPTX_AVAILABLE:
        return docs
    for p in paths:
        p = Path(p)
        if not p.exists():
            continue
        try:
            pres = pptx.Presentation(str(p))
            parts = []
            for i, slide in enumerate(pres.slides):
                parts.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            content = "\n".join([t for t in parts if t and t.strip()])
            docs.append(Doc(content=content, metadata={"source": str(p), "type": "pptx"}))
        except Exception:
            continue
    return docs


def load_xlsx(paths: Iterable[str | Path], max_rows: int = 1000) -> List[Doc]:
    """
    Flatten sheets into a simple text representation (header + first N rows).
    This preserves *some* structure without overcomplicating.
    """
    docs: List[Doc] = []
    if not _XLSX_AVAILABLE:
        return docs
    for p in paths:
        p = Path(p)
        if not p.exists():
            continue
        try:
            xls = pd.ExcelFile(str(p))
            parts = []
            for sheet in xls.sheet_names:
                df = xls.parse(sheet)
                if len(df) > max_rows:
                    df = df.head(max_rows)
                parts.append(f"--- Sheet: {sheet} ---")
                # Convert to a compact, pipe-delimited text block
                parts.append(df.to_csv(index=False, sep="|"))
            content = "\n".join(parts)
            docs.append(Doc(content=content, metadata={"source": str(p), "type": "xlsx"}))
        except Exception:
            continue
    return docs
